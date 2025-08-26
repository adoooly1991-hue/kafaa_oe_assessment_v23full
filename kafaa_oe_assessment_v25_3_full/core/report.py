
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4

def add_header(slide, brand_primary="#C00000", logo_path=None):
    tx = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(0.6))
    p = tx.text_frame.paragraphs[0]; p.text = "Kafaa • OE Assessment"; p.font.size = Pt(18); p.font.bold = True
    if logo_path:
        try: slide.shapes.add_picture(logo_path, Inches(8), Inches(0.05), height=Inches(0.45))
        except: pass

def color(shape, hexcolor):
    r = int(hexcolor[1:3],16); g=int(hexcolor[3:5],16); b=int(hexcolor[5:7],16)
    fill = shape.fill; fill.solid(); fill.fore_color.rgb = RGBColor(r,g,b); return shape

def add_lanes_slide(prs, steps, push_pull=None, stores_flags=None, brand_primary="#C00000", logo_path=None, symbols_dir="assets/symbols"):
    s = prs.slides.add_slide(prs.slide_layouts[5]); add_header(s, brand_primary, logo_path)
    s.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.4)).text_frame.text = "Value Stream — Material & Information Lanes"
    mat = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.1), Inches(9.2), Inches(2.0)); color(mat, "#f7f7f9")
    info = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(3.3), Inches(9.2), Inches(2.0)); color(info, "#f7f7f9")
    s.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(2), Inches(0.3)).text_frame.text = "Material Flow"
    s.shapes.add_textbox(Inches(0.5), Inches(3.25), Inches(2), Inches(0.3)).text_frame.text = "Information Flow"
    steps = steps or []; n = max(1,len(steps)); box_w = 8.6/n; x = 0.9
    for i, stp in enumerate(steps):
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(box_w), Inches(0.9))
        color(box, "#ffffff"); box.line.color.rgb = RGBColor(200,200,200)
        tb = s.shapes.add_textbox(Inches(x+0.06), Inches(1.5), Inches(box_w-0.12), Inches(0.9)); tb.text_frame.text = str(stp)
        if i < len(steps)-1:
            s.shapes.add_connector(1, Inches(x+box_w), Inches(1.95), Inches(x+box_w+0.3), Inches(1.95))
        sig = (push_pull or {}).get(stp, None); flags = (stores_flags or {}).get(stp, {})
        try:
            if sig == "Pull": s.shapes.add_picture(f"{symbols_dir}/pull.png", Inches(x+0.05), Inches(3.6), height=Inches(0.55))
            elif sig == "Push": s.shapes.add_picture(f"{symbols_dir}/push.png", Inches(x+0.05), Inches(3.6), height=Inches(0.55))
            if flags.get("supermarket"): s.shapes.add_picture(f"{symbols_dir}/supermarket.png", Inches(x+0.7), Inches(3.6), height=Inches(0.55))
            if flags.get("fifo"): s.shapes.add_picture(f"{symbols_dir}/fifo.png", Inches(x+1.35), Inches(3.6), height=Inches(0.55))
            if flags.get("pacemaker"): s.shapes.add_picture(f"{symbols_dir}/pacemaker.png", Inches(x+2.0), Inches(3.6), height=Inches(0.55))
        except: pass
        x += box_w + 0.2
    s.shapes.add_textbox(Inches(0.6), Inches(4.9), Inches(8.5), Inches(0.3)).text_frame.text = "Scheduling • Kanban • Production control • Andon"
    try:
        s.shapes.add_textbox(Inches(6.9), Inches(0.9), Inches(2.4), Inches(0.3)).text_frame.text = "Legend"
        s.shapes.add_picture("assets/symbols/pull.png", Inches(6.9), Inches(1.2), height=Inches(0.5)); s.shapes.add_textbox(Inches(7.6), Inches(1.25), Inches(1.0), Inches(0.3)).text_frame.text='Pull'
        s.shapes.add_picture("assets/symbols/push.png", Inches(6.9), Inches(1.8), height=Inches(0.5)); s.shapes.add_textbox(Inches(7.6), Inches(1.85), Inches(1.0), Inches(0.3)).text_frame.text='Push'
        s.shapes.add_picture("assets/symbols/supermarket.png", Inches(8.3), Inches(1.2), height=Inches(0.5)); s.shapes.add_textbox(Inches(9.0), Inches(1.25), Inches(1.4), Inches(0.3)).text_frame.text='Supermarket'
        s.shapes.add_picture("assets/symbols/fifo.png", Inches(8.3), Inches(1.8), height=Inches(0.5)); s.shapes.add_textbox(Inches(9.0), Inches(1.85), Inches(1.0), Inches(0.3)).text_frame.text='FIFO'
    except: pass
    return s

def add_impact_slide(prs, wf_imgs, brand_primary="#C00000", logo_path=None):
    s = prs.slides.add_slide(prs.slide_layouts[5]); add_header(s, brand_primary, logo_path)
    s.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.4)).text_frame.text = "Impact — Savings • Frozen Cash • Sales Opportunity"
    if wf_imgs:
        if wf_imgs.get('savings'): s.shapes.add_picture(wf_imgs['savings'], Inches(0.5), Inches(1.3), width=Inches(3.0))
        if wf_imgs.get('frozen_cash'): s.shapes.add_picture(wf_imgs['frozen_cash'], Inches(3.6), Inches(1.3), width=Inches(3.0))
        if wf_imgs.get('sales_opp'): s.shapes.add_picture(wf_imgs['sales_opp'], Inches(6.7), Inches(1.3), width=Inches(3.0))
    return s

def add_defense_slide(prs, def_kpis=None, compliance=None, brand_primary="#C00000", logo_path=None):
    s = prs.slides.add_slide(prs.slide_layouts[5]); add_header(s, brand_primary, logo_path)
    s.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.4)).text_frame.text = "Defense KPIs & Compliance"
    if def_kpis:
        if def_kpis.get("defense_kpis_radar"): s.shapes.add_picture(def_kpis["defense_kpis_radar"], Inches(0.5), Inches(1.3), width=Inches(4.4))
        if def_kpis.get("defense_kpis_bars"): s.shapes.add_picture(def_kpis["defense_kpis_bars"], Inches(5.0), Inches(1.3), width=Inches(4.4))
    if compliance:
        y = 4.6
        box = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(1.4))
        tf = box.text_frame; tf.word_wrap=True
        tf.text = f"AS9100: {compliance.get('AS9100','-')} • ITAR/EAR: {compliance.get('ITAR/EAR','-')} • FAI%: {compliance.get('FAI%','-')} • Calibration on-time%: {compliance.get('Calibration on-time%','-')}"
        p = tf.add_paragraph(); p.text = f"Trace retrieval (min): {compliance.get('Trace retrieval (min)','-')} • FOD/k hrs: {compliance.get('FOD/k hrs','-')} • Open audit NC: {compliance.get('Open audit NC','-')} • Training%: {compliance.get('Training%','-')}"; p.level=1
    return s

def export_pptx(path, observations=None, actions_df=None, kpi_imgs=None, wc_imgs=None, wf_imgs=None, def_kpis=None, compliance=None, ct_table=None, takt_sec=None, brand_primary="#C00000", logo_path=None, steps=None, push_pull=None, stores_flags=None, evidence=None, review_notes=None):
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5]); add_header(s, brand_primary, logo_path)
    s.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.4)).text_frame.text = "Executive Dashboard"
    if kpi_imgs:
        if kpi_imgs.get('ct_vs_takt'): s.shapes.add_picture(kpi_imgs['ct_vs_takt'], Inches(0.5), Inches(1.3), width=Inches(4.6))
        if kpi_imgs.get('sankey'): s.shapes.add_picture(kpi_imgs['sankey'], Inches(5.3), Inches(1.3), width=Inches(4.6))
    add_lanes_slide(prs, steps or [r.get("step") for r in (ct_table or [])], push_pull=push_pull, stores_flags=stores_flags, brand_primary=brand_primary, logo_path=logo_path)
    if wf_imgs: 
        try: add_impact_slide(prs, wf_imgs, brand_primary, logo_path)
        except: pass
    try: add_defense_slide(prs, def_kpis=kpi_imgs, compliance=compliance, brand_primary=brand_primary, logo_path=logo_path)
    except: pass
    if wc_imgs and (wc_imgs.get('inv_heatmap') or wc_imgs.get('ccc_trend')):
        s = prs.slides.add_slide(prs.slide_layouts[5]); add_header(s, brand_primary, logo_path)
        s.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.4)).text_frame.text = "Working Capital"
        y = 1.3
        if wc_imgs.get('inv_heatmap'): s.shapes.add_picture(wc_imgs['inv_heatmap'], Inches(0.5), Inches(y), width=Inches(9.2)); y += 3.0
        if wc_imgs.get('ccc_trend'): s.shapes.add_picture(wc_imgs['ccc_trend'], Inches(0.5), Inches(y), width=Inches(9.2))
    if observations:
        s = prs.slides.add_slide(prs.slide_layouts[5]); add_header(s, brand_primary, logo_path)
        tb = s.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5)); tf = tb.text_frame; tf.word_wrap=True; tf.text="Observations (PQCDSM)"
        for ptxt in observations: p = tf.add_paragraph(); p.text = ptxt; p.level=1
    if evidence:
        s = prs.slides.add_slide(prs.slide_layouts[5]); add_header(s, brand_primary, logo_path)
        s.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.4)).text_frame.text = "Evidence (Photos)"
        x, y = 0.5, 1.3
        for path in evidence[:8]:
            try:
                s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(3.0))
                x += 3.2
                if x > 9.0: x = 0.5; y += 2.4
            except: pass
    if actions_df is not None and len(actions_df.index)>0:
        s = prs.slides.add_slide(prs.slide_layouts[5]); add_header(s, brand_primary, logo_path)
        tb = s.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5)); tf = tb.text_frame; tf.word_wrap=True; tf.text="Countermeasures (PACE)"
        for _,r in actions_df.iterrows():
            p = tf.add_paragraph(); p.text = f"• [{r.get('PACE_priority','')}] {r.get('stage','')} — {r.get('waste','')}: {r.get('action','')} ({r.get('est_annual_benefit','')})"; p.level=1
    if review_notes:
        s = prs.slides.add_slide(prs.slide_layouts[5]); add_header(s, brand_primary, logo_path)
        s.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.4)).text_frame.text = "Review Notes"
        y = 1.2
        for r in review_notes[:14]:
            tb = s.shapes.add_textbox(Inches(0.6), Inches(y), Inches(8.8), Inches(0.35)); tb.text_frame.text = f"• [{r.get('severity','')}] {r.get('anchor','')}={r.get('ref','')}: {r.get('note','')}"; y += 0.4
    prs.save(path)

def export_pdf(path, paragraphs=None, kpi_imgs=None, wc_imgs=None, wf_imgs=None, logo_path=None):
    c = canvas.Canvas(path, pagesize=A4); w,h = A4
    if logo_path:
        try: c.saveState(); c.translate(w*0.5, h*0.6); c.rotate(30); c.setFillAlpha(0.06); c.drawImage(logo_path, -200, -80, width=400, height=160, mask='auto'); c.restoreState()
        except: pass
    c.setFont("Helvetica-Bold", 18); c.drawString(40, h-40, "Kafaa • OE Assessment")
    c.setFont("Helvetica-Bold", 14); c.drawString(40, h-70, "Executive Dashboard")
    if kpi_imgs:
        if kpi_imgs.get('ct_vs_takt'): c.drawImage(kpi_imgs['ct_vs_takt'], 40, h-320, width=260, height=170, preserveAspectRatio=True, anchor='n')
        if kpi_imgs.get('sankey'): c.drawImage(kpi_imgs['sankey'], 320, h-320, width=260, height=170, preserveAspectRatio=True, anchor='n')
    if wf_imgs and (wf_imgs.get('savings') or wf_imgs.get('frozen_cash') or wf_imgs.get('sales_opp')):
        c.showPage()
        if logo_path:
            try: c.saveState(); c.translate(w*0.5, h*0.6); c.rotate(30); c.setFillAlpha(0.06); c.drawImage(logo_path, -200, -80, width=400, height=160, mask='auto'); c.restoreState()
            except: pass
        c.setFont("Helvetica-Bold", 14); c.drawString(40, h-50, "Impact — Savings • Frozen Cash • Sales Opportunity")
        x = 40
        if wf_imgs.get('savings'): c.drawImage(wf_imgs['savings'], x, h-330, width=170, height=220, preserveAspectRatio=True, anchor='n'); x += 190
        if wf_imgs.get('frozen_cash'): c.drawImage(wf_imgs['frozen_cash'], x, h-330, width=170, height=220, preserveAspectRatio=True, anchor='n'); x += 190
        if wf_imgs.get('sales_opp'): c.drawImage(wf_imgs['sales_opp'], x, h-330, width=170, height=220, preserveAspectRatio=True, anchor='n')
    c.showPage()
    if logo_path:
        try: c.saveState(); c.translate(w*0.5, h*0.6); c.rotate(30); c.setFillAlpha(0.06); c.drawImage(logo_path, -200, -80, width=400, height=160, mask='auto'); c.restoreState()
        except: pass
    c.setFont("Helvetica-Bold", 14); c.drawString(40, h-50, "Observations (PQCDSM)")
    c.setFont("Helvetica", 11); y = h-80
    for p in (paragraphs or []):
        for line in p.split("\\n"):
            c.drawString(40, y, line[:110]); y -= 14
            if y < 60: c.showPage(); y = h-60
    if wc_imgs and (wc_imgs.get('inv_heatmap') or wc_imgs.get('ccc_trend')):
        c.showPage()
        if logo_path:
            try: c.saveState(); c.translate(w*0.5, h*0.6); c.rotate(30); c.setFillAlpha(0.06); c.drawImage(logo_path, -200, -80, width=400, height=160, mask='auto'); c.restoreState()
            except: pass
        c.setFont("Helvetica-Bold", 14); c.drawString(40, h-50, "Working Capital")
        if wc_imgs.get('inv_heatmap'): c.drawImage(wc_imgs['inv_heatmap'], 40, h-240, width=500, height=160, preserveAspectRatio=True, anchor='n')
        if wc_imgs.get('ccc_trend'): c.drawImage(wc_imgs['ccc_trend'], 40, h-430, width=500, height=160, preserveAspectRatio=True, anchor='n')
    c.save()
